'''
PROCESSING SCRIPT FOR TEXT EXTRACTION FROM VARIOUS FILE TYPES

This script extracts text from various file types including .txt, .docx, .doc, .odt, .pptx, .ppt, .xlsx, .xls, .csv, .pdf, and image files (.jpg, .jpeg, .png).
The file types that have been tested are: pdf (with and without OCR), docx, pptx, xlsx, jpg. The rest are theoretically supported but not tested.

Install the dependecies per README.md and check the paths below for Tesseract and Poppler.
'''

#----------------------------------
# LIBRARIES
#----------------------------------

import os
import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd
from pdf2image import convert_from_path
import pytesseract
import cv2
import numpy as np
from PIL import Image
import ollama
from odf import text, teletype
from odf.opendocument import load
import win32com.client  # For .doc, .ppt, .xls on Windows
# import comtypes.client  # Alternative for COM automation

#----------------------------------
# DEPENDENCIES
#----------------------------------

# PATHS OF PACKAGES TO INSTALL - CHECK AND AMEND IF NEEDED
# Set Tesseract path for OCR. 
pytesseract.pytesseract.tesseract_cmd = r'C:/Program Files/Tesseract-OCR/tesseract.exe'
# Set Poppler path for PDF to image processing
actualpoppler_path = r'C:/Users/sanc615/AppData/Local/poppler-24.08.0/Library/bin'

#----------------------------------
# PARAMETERS
#----------------------------------

# OCR mode
ocrmode = True  # Set to True for OCR mode (text from images), False for text-only (faster) PDF processing
# OCR DPI Image resolution
dpino = 150  # Lower DPI for faster processing, adjust as needed
# Max chunk length for LLM processing
maxlength = 3000  # Adjust based on LLM's token limit. Note prompt must be included. Deepseek-R1 has a 4096 token limit.
# Troubleshooting mode
troubleshoot = True  # Set to True to enable txt exports for debugging and troubleshooting

#----------------------------------
# READERS
#----------------------------------

def read_image(file_path):
    """Read text from an image file (.jpg, .jpeg, .png) using Tesseract OCR."""
    try:
        # Verify file exists and is an image
        if not os.path.exists(file_path):
            return f"Error: Image file {file_path} does not exist."
        
        # Load and preprocess image for better OCR accuracy
        img = cv2.imread(file_path, cv2.IMREAD_GRAYSCALE)
        if img is None:
            return f"Error: Failed to load image {file_path}."
        
        # Apply thresholding to improve text clarity
        _, img = cv2.threshold(img, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        
        # Perform OCR
        text = pytesseract.image_to_string(img, lang='eng')
        return text.strip() if text.strip() else "No text extracted from image."
    except Exception as e:
        return f"Error reading image file {file_path}: {str(e)}"

def read_txt(file_path):
    """Read content from a .txt file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        return f"Error reading .txt file: {str(e)}"

def read_docx(file_path):
    """Read content from a .docx file."""
    try:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text])
    except Exception as e:
        return f"Error reading .docx file: {str(e)}"

def read_doc(file_path):
    """Read content from a .doc file (Windows only)."""
    try:
        word = win32com.client.Dispatch("Word.Application")
        doc = word.Documents.Open(os.path.abspath(file_path))
        text = doc.Content.Text
        doc.Close()
        word.Quit()
        return text
    except Exception as e:
        return f"Error reading .doc file: {str(e)}. Note: Requires Microsoft Word on Windows."

def read_odt(file_path):
    """Read content from an .odt file."""
    try:
        doc = load(file_path)
        text_content = ""
        for element in doc.getElementsByType(text.P):
            text_content += teletype.extractText(element) + "\n"
        return text_content
    except Exception as e:
        return f"Error reading .odt file: {str(e)}"

def read_pptx(file_path):
    """Read content from a .pptx file."""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error reading .pptx file: {str(e)}"

def read_ppt(file_path):
    """Read content from a .ppt file (Windows only)."""
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        prs = powerpoint.Presentations.Open(os.path.abspath(file_path))
        text = ""
        for slide in prs.Slides:
            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    text += shape.TextFrame.TextRange.Text + "\n"
        prs.Close()
        powerpoint.Quit()
        return text
    except Exception as e:
        return f"Error reading .ppt file: {str(e)}. Note: Requires Microsoft PowerPoint on Windows."

def read_excel(file_path):
    """Read content from a .xlsx file."""
    try:
        df = pd.read_excel(file_path, engine='openpyxl')
        return df.to_markdown()  # Convert to markdown for better LLM readability
    except Exception as e:
        return f"Error reading .xlsx file: {str(e)}"

def read_xls(file_path):
    """Read content from a .xls file (Windows only)."""
    try:
        excel = win32com.client.Dispatch("Excel.Application")
        wb = excel.Workbooks.Open(os.path.abspath(file_path))
        text = ""
        for sheet in wb.Sheets:
            text += f"Sheet: {sheet.Name}\n"
            text += pd.read_excel(file_path, sheet_name=sheet.Name).to_markdown() + "\n"
        wb.Close()
        excel.Quit()
        return text
    except Exception as e:
        return f"Error reading .xls file: {str(e)}. Note: Requires Microsoft Excel on Windows."

def read_csv(file_path):
    """Read content from a .csv file."""
    try:
        df = pd.read_csv(file_path)
        return df.to_markdown()  # Convert to markdown for better LLM readability
    except Exception as e:
        return f"Error reading .csv file: {str(e)}"

def read_pdf(file_path, ocr_mode=ocrmode):
    """Read content from a .pdf file."""
    try:
        if not ocr_mode:
            with open(file_path, "rb") as file:
                reader = PyPDF2.PdfReader(file)
                text = " ".join(page.extract_text() for page in reader.pages if page.extract_text().strip())
                return text.strip() if text.strip() else "No text extracted from PDF."
        images = convert_from_path(file_path, poppler_path=actualpoppler_path, dpi=dpino)
        text = ""
        for image in images:
            img = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2GRAY)
            img = cv2.threshold(img, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1]
            page_text = pytesseract.image_to_string(img, lang="eng")
            text += page_text.strip() + "\n" if page_text.strip() else ""
        return text.strip() if text.strip() else "No text extracted from PDF via OCR."
    except Exception as e:
        return f"Error reading .pdf file {file_path}: {str(e)}"
    

#----------------------------------
# PROCESSORS
#----------------------------------

# Chunk content into manageable pieces with metadata for LLM processing if files are directly sent to LLMs. The other option
# that we're exploring is a RAG approach with Streamlit and Ollama, which would not require chunking.

# def chunk_content(text, metadata, max_length=maxlength):
#     """Split text into chunks with metadata."""
#     if "Error" in text.lower():
#         return [(text, metadata)]
#     chunks = []
#     num_chunks = (len(text) + max_length - 1) // max_length  # Ceiling division
#     for i in range(0, len(text), max_length):
#         chunk_text = text[i:i + max_length]
#         chunk_metadata = metadata.copy()
#         chunk_metadata["part"] = i // max_length + max_length
#         chunk_metadata["total_parts"] = num_chunks
#         chunks.append((chunk_text.strip(), chunk_metadata))
#     return chunks if chunks else [(f"No text in chunk.", metadata)]


def extract_text_from_file(file_path, ocr_mode=False, troubleshoot=False):
    """Extract text from a single file, with optional troubleshooting."""
    file_extension = os.path.splitext(file_path)[1].lower()
    readers = {
        '.txt': read_txt,
        '.docx': read_docx,
        '.doc': read_doc,
        '.odt': read_odt,
        '.pptx': read_pptx,
        '.ppt': read_ppt,
        '.xlsx': read_excel,
        '.xls': read_xls,
        '.csv': read_csv,
        '.pdf': lambda x: read_pdf(x, ocr_mode=ocr_mode),
        '.jpg': read_image,
        '.jpeg': read_image,
        '.png': read_image
    }

    reader = readers.get(file_extension)
    if not reader:
        return f"Unsupported file type: {file_extension}"
    
    content = reader(file_path)
    if "Error" in content:
        return content

    # # Chunk content
    chunks = chunk_content(content, max_length=maxlength)
    # For simplicity, we can return the full content instead of chunks
    # chunks = [content]  # Treat the entire content as a single chunk for now

    # Troubleshooting: Save chunks to txt_processing subfolder
    if troubleshoot:
        os.makedirs("txt_processing", exist_ok=True)
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        for i, chunk in enumerate(chunks):
            chunk_file = os.path.join("txt_processing", f"{base_name}_chunk_{i}.txt")
            try:
                with open(chunk_file, 'w', encoding='utf-8') as f:
                    f.write(chunk)
            except Exception as e:
                print(f"Warning: Failed to save chunk {chunk_file}: {str(e)}")

    # Return full content (not chunks) for simplicity
    return content

def extract_text_from_directory(directory, ocr_mode=False, troubleshoot=False):
    """Extract text from all supported files in a directory."""
    supported_extensions = {'.txt', '.docx', '.doc', '.odt', '.pptx', '.ppt', '.xlsx', '.xls', '.csv', '.pdf', '.jpg', '.jpeg', '.png'}
    file_paths = []
    results = {}

    # Find supported files
    for root, _, files in os.walk(directory):
        for file in files:
            if os.path.splitext(file)[1].lower() in supported_extensions:
                file_paths.append(os.path.join(root, file))

    if not file_paths:
        return {directory: "No supported files found in directory."}

    print(f"Found {len(file_paths)} supported files: {file_paths}")
    for file_path in file_paths:
        results[file_path] = extract_text_from_file(file_path, ocr_mode, troubleshoot)

    return results

def main():
    """Main function to extract text from files or directories."""
    mode = input("Process a single file or all files in a directory? (single/directory): ").lower()
    ocr_mode = input("Process PDFs with full OCR (captures diagrams, slower)? (yes/no): ").lower() == 'yes'
    troubleshoot = input("Enable troubleshooting (save chunks to txt_processing)? (yes/no): ").lower() == 'yes'

    if mode == 'single':
        file_path = input("Enter the file path: ")
        result = extract_text_from_file(file_path, ocr_mode, troubleshoot)
        print(f"\nResult for {file_path}:\n{result}")
    elif mode == 'directory':
        directory = input("Enter the directory path: ")
        results = extract_text_from_directory(directory, ocr_mode, troubleshoot)
        for file, result in results.items():
            print(f"\n{file}:\n{result}")
    else:
        print("Invalid mode. Choose 'single' or 'directory'.")

if __name__ == "__main__":
    main()